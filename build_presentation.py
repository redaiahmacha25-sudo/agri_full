import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
import zipfile

# Presentation dimensions 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette (Executive Modern Tech Theme)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)       # Slate 50
COLOR_CARD_BG = RGBColor(255, 255, 255)        # Pure White
COLOR_CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200
COLOR_PRIMARY_DARK = RGBColor(15, 23, 42)      # Slate 900
COLOR_PRIMARY_NAVY = RGBColor(30, 41, 59)      # Slate 800
COLOR_ACCENT_BLUE = RGBColor(37, 99, 235)      # Blue 600
COLOR_ACCENT_LIGHT = RGBColor(239, 246, 255)   # Blue 50
COLOR_CYAN = RGBColor(8, 145, 178)             # Cyan 600
COLOR_CYAN_LIGHT = RGBColor(236, 254, 255)     # Cyan 50
COLOR_TEXT_MAIN = RGBColor(15, 23, 42)         # Slate 900
COLOR_TEXT_MUTED = RGBColor(71, 85, 105)       # Slate 600
COLOR_TEXT_LIGHT = RGBColor(148, 163, 184)     # Slate 400
COLOR_WHITE = RGBColor(255, 255, 255)

# Advantage / Disadvantage Colors
COLOR_GREEN_BG = RGBColor(240, 253, 244)       # Green 50
COLOR_GREEN_BORDER = RGBColor(187, 247, 208)   # Green 200
COLOR_GREEN_ACCENT = RGBColor(22, 163, 74)     # Green 600
COLOR_GREEN_TEXT = RGBColor(20, 83, 45)        # Green 900

COLOR_RED_BG = RGBColor(254, 242, 242)         # Red 50
COLOR_RED_BORDER = RGBColor(254, 202, 202)     # Red 200
COLOR_RED_ACCENT = RGBColor(220, 38, 38)       # Red 600
COLOR_RED_TEXT = RGBColor(127, 29, 29)         # Red 900

FONT_NAME = "Calibri"

def set_font(run_or_p, name=FONT_NAME, size=None, bold=None, color=None, italic=None):
    """Helper to set Calibri font properties explicitly"""
    font = run_or_p.font
    font.name = name
    if size is not None:
        font.size = size
    if bold is not None:
        font.bold = bold
    if color is not None:
        font.color.rgb = color
    if italic is not None:
        font.italic = italic

def add_header(slide, category, title, subtitle=None):
    """Adds a standard modern header with a category badge and title"""
    # Category badge background
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.2), Inches(0.32))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_ACCENT_LIGHT
    badge.line.color.rgb = RGBColor(191, 219, 254) # Blue 200
    badge.line.width = Pt(1)
    
    tf = badge.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = category.upper()
    set_font(r, FONT_NAME, Pt(9.5), bold=True, color=COLOR_ACCENT_BLUE)
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.7), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    r_title = p_title.add_run()
    r_title.text = title
    set_font(r_title, FONT_NAME, Pt(22), bold=True, color=COLOR_PRIMARY_DARK)
    
    if subtitle:
        p_sub = tf_title.add_paragraph()
        r_sub = p_sub.add_run()
        r_sub.text = subtitle
        set_font(r_sub, FONT_NAME, Pt(11), bold=False, color=COLOR_TEXT_MUTED)

print("Helper definitions loaded.")
