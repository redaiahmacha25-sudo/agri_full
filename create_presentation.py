import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
import zipfile

def create_presentation(output_path):
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank slide layout (index 6 in default templates, or blank slide)
    blank_layout = prs.slide_layouts[6]
    
    # Color definitions
    C_BG_DARK = RGBColor(15, 23, 42)         # #0F172A Slate 900
    C_BG_DARK_CARD = RGBColor(30, 41, 59)    # #1E293B Slate 800
    C_BG_DARK_CARD_2 = RGBColor(51, 65, 85)  # #334155 Slate 700
    
    C_BG_PAGE = RGBColor(248, 250, 252)      # #F8FAFC Slate 50
    C_CARD_WHITE = RGBColor(255, 255, 255)   # #FFFFFF
    C_CARD_BORDER = RGBColor(226, 232, 240)  # #E2E8F0 Slate 200
    
    C_PRIMARY = RGBColor(15, 23, 42)         # #0F172A
    C_SECONDARY = RGBColor(51, 65, 85)       # #334155
    C_MUTED = RGBColor(100, 116, 139)        # #64748B
    
    C_BLUE = RGBColor(37, 99, 235)           # #2563EB Blue 600
    C_BLUE_LIGHT = RGBColor(239, 246, 255)   # #EFF6FF Blue 50
    C_BLUE_BORDER = RGBColor(191, 219, 254)  # #BFDBFE Blue 200
    C_BLUE_DARK = RGBColor(30, 64, 175)      # #1E40AF Blue 800
    
    C_CYAN = RGBColor(8, 145, 178)           # #0891B2 Cyan 600
    C_CYAN_LIGHT = RGBColor(236, 254, 255)   # #ECFEFF Cyan 50
    C_CYAN_BORDER = RGBColor(165, 243, 252)  # #A5F3FC Cyan 200
    
    C_INDIGO = RGBColor(79, 70, 229)         # #4F46E5 Indigo 600
    C_INDIGO_LIGHT = RGBColor(238, 242, 255) # #EEF2FF Indigo 50
    C_INDIGO_BORDER = RGBColor(199, 210, 254)# #C7D2FE Indigo 200
    
    C_GREEN = RGBColor(22, 163, 74)          # #16A34A Green 600
    C_GREEN_BG = RGBColor(240, 253, 244)     # #F0FDF4 Green 50
    C_GREEN_BORDER = RGBColor(187, 247, 208) # #BBF7D0 Green 200
    C_GREEN_TEXT = RGBColor(20, 83, 45)      # #14532D Green 900
    
    C_RED = RGBColor(220, 38, 38)            # #DC2626 Red 600
    C_RED_BG = RGBColor(254, 242, 242)       # #FEF2F2 Red 50
    C_RED_BORDER = RGBColor(254, 202, 202)   # #FECACA Red 200
    C_RED_TEXT = RGBColor(127, 29, 29)       # #7F1D1D Red 900
    
    FONT_FAMILY = "Calibri"

    def apply_run_font(run, size=None, bold=None, color=None, italic=None):
        run.font.name = FONT_FAMILY
        if size is not None:
            run.font.size = size
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        if italic is not None:
            run.font.italic = italic

    def add_slide_background(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_light_header(slide, category_text, title_text, subtitle_text):
        # Background canvas
        add_slide_background(slide, C_BG_PAGE)
        
        # Category Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.6), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_BLUE_LIGHT
        badge.line.color.rgb = C_BLUE_BORDER
        badge.line.width = Pt(1)
        tf_b = badge.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        r_b = p_b.add_run()
        r_b.text = category_text.upper()
        apply_run_font(r_b, size=Pt(9.5), bold=True, color=C_BLUE)

        # Title & Subtitle box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.76), Inches(11.733), Inches(0.75))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title_text
        apply_run_font(r1, size=Pt(20), bold=True, color=C_PRIMARY)
        
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle_text
        apply_run_font(r2, size=Pt(11), bold=False, color=C_MUTED)

    # =========================================================================
    # SLIDE 1: Title & Agenda Slide
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide1, C_BG_DARK)
    
    # Left Hero Section
    # Category pill
    pill1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.1), Inches(3.2), Inches(0.38))
    pill1.fill.solid()
    pill1.fill.fore_color.rgb = RGBColor(30, 58, 138) # Navy Blue
    pill1.line.color.rgb = RGBColor(96, 165, 250)    # Blue 400
    pill1.line.width = Pt(1)
    tf_p1 = pill1.text_frame
    tf_p1.word_wrap = True
    tf_p1.margin_left = tf_p1.margin_right = tf_p1.margin_top = tf_p1.margin_bottom = 0
    p = tf_p1.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "ROBOTICS & ARTIFICIAL INTELLIGENCE"
    apply_run_font(r, size=Pt(10), bold=True, color=RGBColor(191, 219, 254))
    
    # Main Title Box
    title_box1 = slide1.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(5.6), Inches(3.0))
    tf1 = title_box1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    
    p = tf1.paragraphs[0]
    r = p.add_run()
    r.text = "AI in Robotics:"
    apply_run_font(r, size=Pt(28), bold=True, color=RGBColor(255, 255, 255))
    
    p = tf1.add_paragraph()
    r = p.add_run()
    r.text = "Planning to Move"
    apply_run_font(r, size=Pt(36), bold=True, color=RGBColor(96, 165, 250))
    
    p = tf1.add_paragraph()
    p.space_before = Pt(14)
    r = p.add_run()
    r.text = "Autonomous Navigation, Real-Time Path Optimization & Dynamic Obstacle Avoidance for Next-Generation Intelligent Robots."
    apply_run_font(r, size=Pt(13), color=RGBColor(203, 213, 225))
    
    # Tech Tags Pill Row
    tags = ["Autonomous Systems", "Path Planning", "Real-Time AI", "Obstacle Avoidance"]
    tag_x = Inches(0.9)
    for tag in tags:
        tag_w = Inches(1.3) if len(tag) < 14 else Inches(1.6)
        tp = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tag_x, Inches(5.8), tag_w, Inches(0.36))
        tp.fill.solid()
        tp.fill.fore_color.rgb = C_BG_DARK_CARD
        tp.line.color.rgb = C_BG_DARK_CARD_2
        tp.line.width = Pt(1)
        tf_tag = tp.text_frame
        tf_tag.margin_left = tf_tag.margin_right = tf_tag.margin_top = tf_tag.margin_bottom = 0
        p = tf_tag.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = tag
        apply_run_font(r, size=Pt(9), bold=True, color=RGBColor(148, 163, 184))
        tag_x += tag_w + Inches(0.12)
    
    # Right Section: Agenda Card
    agenda_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.1), Inches(5.5), Inches(5.3))
    agenda_card.fill.solid()
    agenda_card.fill.fore_color.rgb = C_BG_DARK_CARD
    agenda_card.line.color.rgb = C_BG_DARK_CARD_2
    agenda_card.line.width = Pt(1)
    
    # Agenda Title Box
    ag_title_box = slide1.shapes.add_textbox(Inches(7.2), Inches(1.35), Inches(4.9), Inches(0.5))
    tf_at = ag_title_box.text_frame
    tf_at.margin_left = tf_at.margin_right = tf_at.margin_top = tf_at.margin_bottom = 0
    p = tf_at.paragraphs[0]
    r = p.add_run()
    r.text = "PRESENTATION AGENDA"
    apply_run_font(r, size=Pt(12), bold=True, color=RGBColor(148, 163, 184))
    
    # 4 Agenda Item Rows
    agenda_items = [
        ("01", "Introduction & Core Concepts", "Definition, perception sensors, and the fundamental navigation pipeline"),
        ("02", "Working Process Lifecycle", "The 8-step execution loop from goal acquisition to safe arrival"),
        ("03", "Applications & Popular Algorithms", "Real-world robotic use cases and key AI algorithms (A*, Dijkstra, D*, RRT, RL)"),
        ("04", "Advantages & Disadvantages", "Comprehensive analysis of capabilities, operational gains, and technical limits")
    ]
    
    ag_y = Inches(1.95)
    for num, title, desc in agenda_items:
        # Number Pill
        np = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), ag_y, Inches(0.65), Inches(0.9))
        np.fill.solid()
        np.fill.fore_color.rgb = RGBColor(30, 58, 138)
        np.line.color.rgb = RGBColor(59, 130, 246)
        np.line.width = Pt(1)
        tf_np = np.text_frame
        tf_np.margin_left = tf_np.margin_right = tf_np.margin_top = tf_np.margin_bottom = 0
        p = tf_np.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        apply_run_font(r, size=Pt(14), bold=True, color=RGBColor(255, 255, 255))
        
        # Item Text Box
        it_box = slide1.shapes.add_textbox(Inches(8.0), ag_y, Inches(4.15), Inches(0.9))
        tf_it = it_box.text_frame
        tf_it.word_wrap = True
        tf_it.margin_left = tf_it.margin_right = tf_it.margin_top = tf_it.margin_bottom = 0
        p = tf_it.paragraphs[0]
        r = p.add_run()
        r.text = title
        apply_run_font(r, size=Pt(12), bold=True, color=RGBColor(255, 255, 255))
        
        p = tf_it.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = desc
        apply_run_font(r, size=Pt(10), color=RGBColor(148, 163, 184))
        
        ag_y += Inches(1.1)

    # =========================================================================
    # SLIDE 2: Planning to Move (Introduction & Overview)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_light_header(slide2, "1. Introduction & Overview", "Planning to Move: Autonomous Robot Navigation", "Core definitions, sensing architecture, and closed-loop real-time obstacle avoidance.")
    
    # 3 Top Concept Cards
    cards_s2 = [
        ("WHAT IS PLANNING TO MOVE?", 
         "The ability of an AI robot to autonomously determine the safest and most efficient path from its current position to a target destination without human intervention.",
         C_BLUE_LIGHT, C_BLUE_BORDER, C_BLUE),
        ("PERCEPTION & ANALYSIS", 
         "The robot continuously gathers telemetry via multi-modal sensors (LiDAR, Cameras, GPS, Ultrasonic) and analyzes environmental geometry to identify static and dynamic obstacles.",
         C_CYAN_LIGHT, C_CYAN_BORDER, C_CYAN),
        ("REAL-TIME REPLANNING", 
         "When an obstacle or sudden environmental shift occurs, onboard AI models dynamically recalculate an alternate collision-free trajectory instantaneously in real time.",
         C_INDIGO_LIGHT, C_INDIGO_BORDER, C_INDIGO)
    ]
    
    card_w = Inches(3.72)
    card_h = Inches(2.25)
    for idx, (head, body, bg_c, border_c, accent_c) in enumerate(cards_s2):
        cx = Inches(0.8 + idx * 3.98)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.7), card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_WHITE
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1)
        
        # Header strip inside card
        strip = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx + Inches(0.2), Inches(1.9), Inches(3.32), Inches(0.35))
        strip.fill.solid()
        strip.fill.fore_color.rgb = bg_c
        strip.line.color.rgb = border_c
        strip.line.width = Pt(1)
        tf_s = strip.text_frame
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = 0
        p = tf_s.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        apply_run_font(r, size=Pt(9.5), bold=True, color=accent_c)
        
        # Body text
        tb = slide2.shapes.add_textbox(cx + Inches(0.2), Inches(2.35), Inches(3.32), Inches(1.45))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = body
        apply_run_font(r, size=Pt(10.5), color=C_SECONDARY)

    # Bottom Pipeline Card
    pipe_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.15), Inches(11.733), Inches(2.75))
    pipe_card.fill.solid()
    pipe_card.fill.fore_color.rgb = C_CARD_WHITE
    pipe_card.line.color.rgb = C_CARD_BORDER
    pipe_card.line.width = Pt(1)
    
    # Pipeline Header
    tb_ph = slide2.shapes.add_textbox(Inches(1.05), Inches(4.3), Inches(11.2), Inches(0.35))
    tf_ph = tb_ph.text_frame
    tf_ph.margin_left = tf_ph.margin_right = tf_ph.margin_top = tf_ph.margin_bottom = 0
    p = tf_ph.paragraphs[0]
    r = p.add_run()
    r.text = "INTELLIGENT AUTONOMOUS NAVIGATION PIPELINE"
    apply_run_font(r, size=Pt(11), bold=True, color=C_PRIMARY)
    
    # 7 Connected Pipeline Steps
    pipeline_steps = [
        ("Environment", "Real-world context", C_BG_PAGE, C_CARD_BORDER, C_PRIMARY),
        ("Sensors", "Cameras, LiDAR, GPS", C_BLUE_LIGHT, C_BLUE_BORDER, C_BLUE),
        ("AI Path Planning", "Route calculation", C_CYAN_LIGHT, C_CYAN_BORDER, C_CYAN),
        ("Obstacle Detection", "Collision sensing", C_INDIGO_LIGHT, C_INDIGO_BORDER, C_INDIGO),
        ("Replanning", "Dynamic rerouting", RGBColor(254, 243, 199), RGBColor(253, 230, 138), RGBColor(180, 83, 9)),
        ("Movement", "Motor control", C_BLUE_LIGHT, C_BLUE_BORDER, C_BLUE),
        ("Destination", "Goal reached", C_GREEN_BG, C_GREEN_BORDER, C_GREEN)
    ]
    
    step_w = Inches(1.38)
    step_h = Inches(1.1)
    for idx, (s_title, s_sub, s_bg, s_bc, s_tc) in enumerate(pipeline_steps):
        sx = Inches(1.05 + idx * 1.62)
        sh = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, Inches(4.75), step_w, step_h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = s_bg
        sh.line.color.rgb = s_bc
        sh.line.width = Pt(1)
        
        tf_s = sh.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_right = tf_s.margin_top = tf_s.margin_bottom = Inches(0.08)
        
        p = tf_s.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{idx+1}\n"
        apply_run_font(r, size=Pt(9), bold=True, color=s_tc)
        
        r2 = p.add_run()
        r2.text = s_title + "\n"
        apply_run_font(r2, size=Pt(9.5), bold=True, color=C_PRIMARY)
        
        r3 = p.add_run()
        r3.text = s_sub
        apply_run_font(r3, size=Pt(8), color=C_MUTED)
        
        # Arrow connecting to next step (except last)
        if idx < len(pipeline_steps) - 1:
            arr_box = slide2.shapes.add_textbox(sx + step_w, Inches(5.1), Inches(0.24), Inches(0.4))
            tf_arr = arr_box.text_frame
            tf_arr.margin_left = tf_arr.margin_right = tf_arr.margin_top = tf_arr.margin_bottom = 0
            p_arr = tf_arr.paragraphs[0]
            p_arr.alignment = PP_ALIGN.CENTER
            r_arr = p_arr.add_run()
            r_arr.text = "➔"
            apply_run_font(r_arr, size=Pt(12), bold=True, color=C_MUTED)

    # Footnote summary
    fn_box = slide2.shapes.add_textbox(Inches(1.05), Inches(6.1), Inches(11.2), Inches(0.5))
    tf_fn = fn_box.text_frame
    tf_fn.word_wrap = True
    tf_fn.margin_left = tf_fn.margin_right = tf_fn.margin_top = tf_fn.margin_bottom = 0
    p_fn = tf_fn.paragraphs[0]
    r_fn = p_fn.add_run()
    r_fn.text = "Key Applications: "
    apply_run_font(r_fn, size=Pt(9.5), bold=True, color=C_PRIMARY)
    r_fn2 = p_fn.add_run()
    r_fn2.text = "Autonomous Vehicles, Warehouse AMRs, Sidewalk Delivery Robots, Medical & Hospital Bots, and Industrial Automation Systems."
    apply_run_font(r_fn2, size=Pt(9.5), color=C_MUTED)

    # =========================================================================
    # SLIDE 3: Planning to Move – Working Process (8 Detailed Steps)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_light_header(slide3, "2. Working Process", "Planning to Move: The 8-Step Navigation Lifecycle", "Step-by-step end-to-end execution loop for autonomous robotic decision-making.")
    
    steps_s3 = [
        ("1", "Receives the Goal", "The robot acquires target destination coordinates from user input, fleet management, or high-level mission directives."),
        ("2", "Senses Environment", "High-frequency cameras, LiDAR scanners, ultrasonic sensors, and GPS gather real-time spatial data around the robot."),
        ("3", "Processes the Data", "AI perception pipelines process point clouds and images to construct 2D/3D costmaps, identify obstacles, and map terrain."),
        ("4", "Plans the Best Path", "Global path planning algorithms evaluate potential trajectories to select the shortest, safest, and most energy-efficient route."),
        ("5", "Executes Movement", "Motor controllers and low-level actuators precisely drive the robot along the calculated path while regulating speed and torque."),
        ("6", "Detects Obstacles", "Proximity sensors continuously scan ahead in real time to detect dynamic obstacles like pedestrians, vehicles, or barriers."),
        ("7", "Replans the Route", "If a path blockage occurs, local replanning algorithms instantly calculate an evasive, collision-free alternative route."),
        ("8", "Reaches Destination", "The robot safely arrives at target coordinates, stabilizes position, and completes the assigned delivery or operation.")
    ]
    
    # 4 columns x 2 rows
    c_w = Inches(2.78)
    c_h = Inches(2.05)
    for idx, (num, s_title, s_desc) in enumerate(steps_s3):
        col = idx % 4
        row = idx // 4
        x = Inches(0.8 + col * 2.98)
        y = Inches(1.68 + row * 2.22)
        
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, c_w, c_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_WHITE
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1)
        
        # Step Number Badge
        nb = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.18), y + Inches(0.18), Inches(0.42), Inches(0.42))
        nb.fill.solid()
        nb.fill.fore_color.rgb = C_BLUE_LIGHT
        nb.line.color.rgb = C_BLUE_BORDER
        nb.line.width = Pt(1)
        tf_nb = nb.text_frame
        tf_nb.margin_left = tf_nb.margin_right = tf_nb.margin_top = tf_nb.margin_bottom = 0
        p = tf_nb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        apply_run_font(r, size=Pt(11), bold=True, color=C_BLUE)
        
        # Step Title
        tb_t = slide3.shapes.add_textbox(x + Inches(0.68), y + Inches(0.18), Inches(1.95), Inches(0.42))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_right = tf_t.margin_top = tf_t.margin_bottom = 0
        p = tf_t.paragraphs[0]
        r = p.add_run()
        r.text = s_title
        apply_run_font(r, size=Pt(11), bold=True, color=C_PRIMARY)
        
        # Step Description
        tb_d = slide3.shapes.add_textbox(x + Inches(0.18), y + Inches(0.68), Inches(2.42), Inches(1.25))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_right = tf_d.margin_top = tf_d.margin_bottom = 0
        p = tf_d.paragraphs[0]
        r = p.add_run()
        r.text = s_desc
        apply_run_font(r, size=Pt(9.5), color=C_SECONDARY)

    # Bottom workflow diagram pill
    wf_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.25), Inches(11.733), Inches(0.75))
    wf_box.fill.solid()
    wf_box.fill.fore_color.rgb = C_BG_DARK_CARD
    wf_box.line.color.rgb = C_BG_DARK_CARD_2
    wf_box.line.width = Pt(1)
    
    tf_wf = wf_box.text_frame
    tf_wf.word_wrap = True
    tf_wf.margin_left = tf_wf.margin_right = tf_wf.margin_top = tf_wf.margin_bottom = 0
    p_wf = tf_wf.paragraphs[0]
    p_wf.alignment = PP_ALIGN.CENTER
    
    wf_nodes = ["1. Goal", "2. Sensing", "3. AI Processing", "4. Path Planning", "5. Movement", "6. Obstacle Detect", "7. Replanning", "8. Destination"]
    for i, node in enumerate(wf_nodes):
        r_n = p_wf.add_run()
        r_n.text = node
        apply_run_font(r_n, size=Pt(10), bold=True, color=RGBColor(255, 255, 255) if i in [0, 3, 7] else RGBColor(147, 197, 253))
        if i < len(wf_nodes) - 1:
            r_sep = p_wf.add_run()
            r_sep.text = "  ➔  "
            apply_run_font(r_sep, size=Pt(10), bold=True, color=RGBColor(148, 163, 184))

    # =========================================================================
    # SLIDE 4: Applications & Popular Algorithms
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_light_header(slide4, "3. Applications & Algorithms", "Robotic Applications & Key AI Path Planning Algorithms", "Bridging practical industry use cases with mathematical navigation and graph search algorithms.")
    
    # Two Columns: Applications (Left) and Algorithms (Right)
    col_w = Inches(5.72)
    col_h = Inches(5.35)
    
    # Left Card: Applications
    card_app = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.65), col_w, col_h)
    card_app.fill.solid()
    card_app.fill.fore_color.rgb = C_CARD_WHITE
    card_app.line.color.rgb = C_CARD_BORDER
    card_app.line.width = Pt(1)
    
    # App Header Banner
    app_banner = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(1.85), Inches(5.22), Inches(0.45))
    app_banner.fill.solid()
    app_banner.fill.fore_color.rgb = C_BLUE_LIGHT
    app_banner.line.color.rgb = C_BLUE_BORDER
    app_banner.line.width = Pt(1)
    tf_ab = app_banner.text_frame
    tf_ab.margin_left = tf_ab.margin_right = tf_ab.margin_top = tf_ab.margin_bottom = 0
    p = tf_ab.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "KEY ROBOTICS APPLICATIONS"
    apply_run_font(r, size=Pt(11), bold=True, color=C_BLUE_DARK)
    
    apps_list = [
        ("1. Autonomous Vehicles", "Navigate complex city traffic, multi-lane highways, and intersections with continuous 360° obstacle avoidance."),
        ("2. Warehouse Robots (AMRs)", "Transport pallets, bins, and cargo swiftly between storage racks and loading docks in high-density facilities."),
        ("3. Sidewalk Delivery Robots", "Perform last-mile delivery of food, groceries, and parcels on pedestrian sidewalks and campuses."),
        ("4. Hospital Service Robots", "Safely transport medicines, lab specimens, linens, and sterile supplies without disrupting doctors or patients."),
        ("5. Drones & Aerial UAVs", "Execute autonomous flight paths for aerial surveillance, infrastructure inspection, 3D mapping, and package drop-off.")
    ]
    
    app_y = Inches(2.45)
    for title, desc in apps_list:
        row_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), app_y, Inches(5.22), Inches(0.8))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = C_BG_PAGE
        row_box.line.color.rgb = C_CARD_BORDER
        row_box.line.width = Pt(1)
        
        tb = slide4.shapes.add_textbox(Inches(1.2), app_y + Inches(0.08), Inches(4.92), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        apply_run_font(r, size=Pt(10.5), bold=True, color=C_PRIMARY)
        
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = desc
        apply_run_font(r, size=Pt(9), color=C_MUTED)
        
        app_y += Inches(0.88)

    # Right Card: Popular Algorithms
    card_algo = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.65), col_w, col_h)
    card_algo.fill.solid()
    card_algo.fill.fore_color.rgb = C_CARD_WHITE
    card_algo.line.color.rgb = C_CARD_BORDER
    card_algo.line.width = Pt(1)
    
    # Algo Header Banner
    algo_banner = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.85), Inches(5.22), Inches(0.45))
    algo_banner.fill.solid()
    algo_banner.fill.fore_color.rgb = C_CYAN_LIGHT
    algo_banner.line.color.rgb = C_CYAN_BORDER
    algo_banner.line.width = Pt(1)
    tf_alb = algo_banner.text_frame
    tf_alb.margin_left = tf_alb.margin_right = tf_alb.margin_top = tf_alb.margin_bottom = 0
    p = tf_alb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "POPULAR AI PATH PLANNING ALGORITHMS"
    apply_run_font(r, size=Pt(11), bold=True, color=C_CYAN)
    
    algos_list = [
        ("1. A* (A-Star) Algorithm", "The benchmark heuristic-guided search algorithm that calculates the shortest and most optimal path with exceptional efficiency."),
        ("2. Dijkstra's Algorithm", "Classic graph search algorithm determining minimum-cost paths in fully mapped, known static environments."),
        ("3. D* Lite (Dynamic A*)", "Optimized incremental heuristic search specifically engineered for real-time dynamic replanning when obstacles appear."),
        ("4. Rapidly-exploring Random Tree (RRT)", "Probabilistic tree-sampling method ideal for high-dimensional configuration spaces and non-holonomic robots."),
        ("5. Reinforcement Learning (RL)", "Neural network policies trained via deep Q-learning / PPO that learn adaptive obstacle avoidance through trial and reward.")
    ]
    
    algo_y = Inches(2.45)
    for title, desc in algos_list:
        row_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), algo_y, Inches(5.22), Inches(0.8))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = C_BG_PAGE
        row_box.line.color.rgb = C_CARD_BORDER
        row_box.line.width = Pt(1)
        
        tb = slide4.shapes.add_textbox(Inches(7.2), algo_y + Inches(0.08), Inches(4.92), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        apply_run_font(r, size=Pt(10.5), bold=True, color=C_PRIMARY)
        
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = desc
        apply_run_font(r, size=Pt(9), color=C_MUTED)
        
        algo_y += Inches(0.88)

    # =========================================================================
    # SLIDE 5: Advantages & Disadvantages
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_light_header(slide5, "4. Advantages & Disadvantages", "Advantages & Disadvantages of AI Path Planning", "Evaluating operational autonomy and navigation efficiency against hardware, sensor, and cost tradeoffs.")
    
    # Left Column: Advantages (Green Theme)
    card_adv = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.65), col_w, col_h)
    card_adv.fill.solid()
    card_adv.fill.fore_color.rgb = C_CARD_WHITE
    card_adv.line.color.rgb = C_GREEN_BORDER
    card_adv.line.width = Pt(1)
    
    adv_banner = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(1.85), Inches(5.22), Inches(0.42))
    adv_banner.fill.solid()
    adv_banner.fill.fore_color.rgb = C_GREEN_BG
    adv_banner.line.color.rgb = C_GREEN_BORDER
    adv_banner.line.width = Pt(1)
    tf_advb = adv_banner.text_frame
    tf_advb.margin_left = tf_advb.margin_right = tf_advb.margin_top = tf_advb.margin_bottom = 0
    p = tf_advb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "KEY ADVANTAGES & BENEFITS"
    apply_run_font(r, size=Pt(11), bold=True, color=C_GREEN_TEXT)
    
    adv_items = [
        ("Efficient Navigation", "Finds the mathematically optimal shortest and safest path to minimize travel time."),
        ("Obstacle Avoidance", "Detects and actively maneuvers around static and dynamic obstacles automatically."),
        ("Saves Time & Energy", "Optimizes motor acceleration and route length, significantly reducing battery consumption."),
        ("Improved Accuracy", "Increases positioning precision, repeatability, and task execution consistency."),
        ("Autonomous Operation", "Operates independently 24/7 with minimal or zero human supervision."),
        ("Adaptable & Flexible", "Instantly recalculates routes in real time when environment conditions change.")
    ]
    
    adv_y = Inches(2.4)
    for title, desc in adv_items:
        row_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), adv_y, Inches(5.22), Inches(0.7))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = C_GREEN_BG
        row_box.line.color.rgb = C_GREEN_BORDER
        row_box.line.width = Pt(1)
        
        tb = slide5.shapes.add_textbox(Inches(1.2), adv_y + Inches(0.06), Inches(4.92), Inches(0.58))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "✔  " + title + ": "
        apply_run_font(r, size=Pt(10), bold=True, color=C_GREEN_TEXT)
        r2 = p.add_run()
        r2.text = desc
        apply_run_font(r2, size=Pt(9), color=C_SECONDARY)
        
        adv_y += Inches(0.74)

    # Right Column: Disadvantages (Red Theme)
    card_dis = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.65), col_w, col_h)
    card_dis.fill.solid()
    card_dis.fill.fore_color.rgb = C_CARD_WHITE
    card_dis.line.color.rgb = C_RED_BORDER
    card_dis.line.width = Pt(1)
    
    dis_banner = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.85), Inches(5.22), Inches(0.42))
    dis_banner.fill.solid()
    dis_banner.fill.fore_color.rgb = C_RED_BG
    dis_banner.line.color.rgb = C_RED_BORDER
    dis_banner.line.width = Pt(1)
    tf_disb = dis_banner.text_frame
    tf_disb.margin_left = tf_disb.margin_right = tf_disb.margin_top = tf_disb.margin_bottom = 0
    p = tf_disb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "DISADVANTAGES & CHALLENGES"
    apply_run_font(r, size=Pt(11), bold=True, color=C_RED_TEXT)
    
    dis_items = [
        ("High Computational Power", "Requires high-performance onboard GPUs, TPUs, and multi-core processors."),
        ("High Implementation Cost", "Specialized LiDAR scanners, cameras, and robotics software development are costly."),
        ("Dependence on Sensors", "Reliability hinges on clean sensor streams; susceptible to blinding light, fog, or dust."),
        ("Complex Environments", "Planning latency and search space explode in dense, highly chaotic crowds."),
        ("Maintenance Required", "Hardware sensors and SLAM calibration parameters demand continuous maintenance."),
        ("Error & Failure Risk", "Sensor occlusion, software edge-cases, or blind spots can cause navigation errors.")
    ]
    
    dis_y = Inches(2.4)
    for title, desc in dis_items:
        row_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), dis_y, Inches(5.22), Inches(0.7))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = C_RED_BG
        row_box.line.color.rgb = C_RED_BORDER
        row_box.line.width = Pt(1)
        
        tb = slide5.shapes.add_textbox(Inches(7.2), dis_y + Inches(0.06), Inches(4.92), Inches(0.58))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "✖  " + title + ": "
        apply_run_font(r, size=Pt(10), bold=True, color=C_RED_TEXT)
        r2 = p.add_run()
        r2.text = desc
        apply_run_font(r2, size=Pt(9), color=C_SECONDARY)
        
        dis_y += Inches(0.74)

    # =========================================================================
    # SLIDE 6: Concluding / Thank You Slide
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide6, C_BG_DARK)
    
    # Outer Centered Glow Card
    thx_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(1.0), Inches(9.733), Inches(5.5))
    thx_card.fill.solid()
    thx_card.fill.fore_color.rgb = C_BG_DARK_CARD
    thx_card.line.color.rgb = RGBColor(59, 130, 246)
    thx_card.line.width = Pt(1.5)
    
    # Topic Badge
    t_badge = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.4), Inches(3.733), Inches(0.38))
    t_badge.fill.solid()
    t_badge.fill.fore_color.rgb = RGBColor(30, 58, 138)
    t_badge.line.color.rgb = RGBColor(96, 165, 250)
    t_badge.line.width = Pt(1)
    tf_tb = t_badge.text_frame
    tf_tb.margin_left = tf_tb.margin_right = tf_tb.margin_top = tf_tb.margin_bottom = 0
    p = tf_tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "AI IN ROBOTICS — PLANNING TO MOVE"
    apply_run_font(r, size=Pt(10), bold=True, color=RGBColor(191, 219, 254))
    
    # Big Thank You Text
    thx_tb = slide6.shapes.add_textbox(Inches(2.2), Inches(2.0), Inches(8.933), Inches(1.1))
    tf_thx = thx_tb.text_frame
    tf_thx.margin_left = tf_thx.margin_right = tf_thx.margin_top = tf_thx.margin_bottom = 0
    p = tf_thx.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "THANK YOU!"
    apply_run_font(r, size=Pt(44), bold=True, color=RGBColor(255, 255, 255))
    
    p = tf_thx.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Questions & Discussion"
    apply_run_font(r, size=Pt(18), bold=True, color=RGBColor(96, 165, 250))
    
    # Key Takeaways Box
    summary_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.4), Inches(3.4), Inches(8.533), Inches(1.7))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = C_BG_DARK
    summary_box.line.color.rgb = C_BG_DARK_CARD_2
    summary_box.line.width = Pt(1)
    
    tf_sum = summary_box.text_frame
    tf_sum.word_wrap = True
    tf_sum.margin_left = Inches(0.3)
    tf_sum.margin_right = Inches(0.3)
    tf_sum.margin_top = Inches(0.2)
    tf_sum.margin_bottom = Inches(0.2)
    
    p = tf_sum.paragraphs[0]
    r = p.add_run()
    r.text = "Key Takeaway Summary:\n"
    apply_run_font(r, size=Pt(11.5), bold=True, color=RGBColor(255, 255, 255))
    
    p2 = tf_sum.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = "Planning to Move is the core intelligence underpinning autonomous robotics. By harmonizing sensor perception (LiDAR, Cameras, GPS) with advanced search algorithms (A*, D* Lite, RRT) and real-time replanning, modern robots achieve safe, efficient, and collision-free mobility across factories, hospitals, roads, and cities."
    apply_run_font(r2, size=Pt(10), color=RGBColor(203, 213, 225))
    
    # Footer info
    ft_tb = slide6.shapes.add_textbox(Inches(2.4), Inches(5.3), Inches(8.533), Inches(0.5))
    tf_ft = ft_tb.text_frame
    tf_ft.margin_left = tf_ft.margin_right = tf_ft.margin_top = tf_ft.margin_bottom = 0
    p = tf_ft.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Autonomous Systems  •  AI Motion Planning  •  Real-Time Collision Avoidance"
    apply_run_font(r, size=Pt(9.5), bold=True, color=RGBColor(148, 163, 184))
    
    # Save the presentation
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == "__main__":
    out = r"C:\Users\Redaiah\OneDrive\Documents\Planning to move.pptx"
    create_presentation(out)
